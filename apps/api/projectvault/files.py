from __future__ import annotations

import csv
import hashlib
import json
import mimetypes
import shutil
import tempfile
import uuid
import zipfile
from datetime import datetime, timezone
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .db import Database
from .ingest import sha256_file, utc_now

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".jsonl", ".csv", ".tsv", ".xml", ".html", ".htm",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".go", ".rs", ".c", ".h", ".cpp", ".hpp",
    ".cs", ".php", ".rb", ".swift", ".kt", ".kts", ".scala", ".sh", ".bash", ".zsh", ".fish",
    ".sql", ".toml", ".yaml", ".yml", ".ini", ".cfg", ".conf", ".env", ".log", ".tex",
}
MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_CELL_COUNT = 200_000
MAX_FAILURE_SAMPLES = 200


@dataclass(frozen=True)
class ArchiveLimits:
    max_files: int
    max_expanded_bytes: int
    max_member_bytes: int
    max_indexable_file_bytes: int
    min_free_bytes: int

    def __post_init__(self) -> None:
        values = {
            "max_files": self.max_files,
            "max_expanded_bytes": self.max_expanded_bytes,
            "max_member_bytes": self.max_member_bytes,
            "max_indexable_file_bytes": self.max_indexable_file_bytes,
            "min_free_bytes": self.min_free_bytes,
        }
        if any(value < 1 for value in values.values()):
            raise ValueError("archive_limits_must_be_positive")
        if self.max_member_bytes > self.max_expanded_bytes:
            raise ValueError("archive_member_limit_exceeds_expanded_limit")
        if self.max_indexable_file_bytes > self.max_member_bytes:
            raise ValueError("indexable_file_limit_exceeds_member_limit")

    @classmethod
    def from_settings(cls, settings: Any) -> "ArchiveLimits":
        return cls(
            max_files=settings.max_archive_files,
            max_expanded_bytes=settings.max_archive_expanded_bytes,
            max_member_bytes=settings.max_archive_member_bytes,
            max_indexable_file_bytes=settings.max_indexable_file_bytes,
            min_free_bytes=settings.min_free_bytes,
        )


def _normalised_archive_path(member_name: str) -> str:
    normalised = member_name.replace("\\", "/")
    if not normalised or normalised.startswith("/") or "\x00" in normalised:
        raise ValueError("unsafe_archive_path")
    parts = [part for part in normalised.split("/") if part not in ("", ".")]
    if not parts or any(part == ".." for part in parts):
        raise ValueError("unsafe_archive_path")
    return "/".join(parts)


def _safe_archive_members(archive: zipfile.ZipFile, limits: ArchiveLimits) -> list[tuple[zipfile.ZipInfo, str]]:
    total = 0
    paths: set[str] = set()
    members: list[tuple[zipfile.ZipInfo, str]] = []
    for info in archive.infolist():
        unix_mode = (info.external_attr >> 16) & 0xFFFF
        if unix_mode and (unix_mode & 0o170000) == 0o120000:
            raise ValueError("archive_symlink_rejected")
        if info.is_dir():
            continue
        if info.flag_bits & 0x1:
            raise ValueError("encrypted_archive_member_rejected")
        relative = _normalised_archive_path(info.filename)
        if relative in paths:
            raise ValueError("duplicate_archive_member")
        paths.add(relative)
        if info.file_size > limits.max_member_bytes:
            raise ValueError("archive_member_size_limit")
        total += info.file_size
        if len(members) + 1 > limits.max_files:
            raise ValueError("archive_file_count_limit")
        if total > limits.max_expanded_bytes:
            raise ValueError("archive_expanded_size_limit")
        members.append((info, relative))
    return members


def _require_space(directory: Path, incoming_bytes: int, min_free_bytes: int) -> None:
    if shutil.disk_usage(directory).free < incoming_bytes + min_free_bytes:
        raise ValueError("insufficient_storage")


def _write_archive_member(
    archive: zipfile.ZipFile,
    info: zipfile.ZipInfo,
    destination: Path,
    limits: ArchiveLimits,
) -> None:
    written = 0
    with archive.open(info) as source, destination.open("wb") as target:
        while chunk := source.read(1024 * 1024):
            written += len(chunk)
            if written > limits.max_member_bytes:
                raise ValueError("archive_member_size_limit")
            _require_space(destination.parent, len(chunk), limits.min_free_bytes)
            target.write(chunk)
    if written != info.file_size:
        raise zipfile.BadZipFile("archive_member_size_mismatch")


def _text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        parts.append(f"## Page {index}\n\n{page.extract_text() or ''}")
    return "\n\n".join(parts)


def _docx(path: Path) -> str:
    document = DocxDocument(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"\n## Table {table_index}")
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    cells = 0
    try:
        for worksheet in workbook.worksheets:
            parts.append(f"# Sheet: {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    parts.append("\t".join(values))
                cells += len(values)
                if cells > MAX_CELL_COUNT:
                    parts.append("[TRUNCATED: spreadsheet cell limit reached]")
                    return "\n".join(parts)
    finally:
        workbook.close()
    return "\n".join(parts)


def _pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                parts.append(text)
    return "\n".join(parts)


def extractor(path: Path) -> Callable[[Path], str] | None:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        return _text
    if suffix == ".pdf":
        return _pdf
    if suffix == ".docx":
        return _docx
    if suffix in {".xlsx", ".xlsm"}:
        return _xlsx
    if suffix == ".pptx":
        return _pptx
    return None


def ingest_directory(
    db: Database,
    root: Path,
    project_id: str,
    project_name: str,
    *,
    max_file_bytes: int = MAX_FILE_BYTES,
) -> dict[str, Any]:
    root = root.resolve()
    if not root.is_dir():
        raise ValueError(f"Not a directory: {root}")
    project_id = project_id.strip()
    project_name = project_name.strip()
    if not project_id or not project_name:
        raise ValueError("project_id and project_name are required")

    started = utc_now()
    db.upsert_project(project_id, project_name, "owner_supplied_directory")
    imported = skipped = failed = 0
    failures: list[dict[str, str]] = []
    aggregate = hashlib.sha256()

    for path in sorted(item for item in root.rglob("*") if item.is_file()):
        relative = path.relative_to(root).as_posix()
        if path.stat().st_size > max_file_bytes:
            skipped += 1
            failures.append({"path": relative, "reason": "file_too_large"})
            continue
        read = extractor(path)
        if read is None:
            skipped += 1
            failures.append({"path": relative, "reason": "unsupported_type"})
            continue
        try:
            text = read(path).strip()
            if not text:
                skipped += 1
                failures.append({"path": relative, "reason": "empty_extracted_text"})
                continue
            content_hash = sha256_file(path)
            aggregate.update(relative.encode("utf-8"))
            aggregate.update(content_hash.encode("ascii"))
            document_id = f"file:{hashlib.sha256((project_id + '\0' + relative).encode('utf-8')).hexdigest()}"
            timestamp = path.stat().st_mtime
            body = (
                f"# {path.name}\n\n## Provenance\n\nProject: {project_name}\n"
                f"Project ID: {project_id}\nAttribution basis: owner_supplied_directory\n"
                f"Relative path: {relative}\nSHA-256: {content_hash}\n\n## Extracted content\n\n{text}\n"
            )
            db.upsert_document({
                "document_id": document_id,
                "conversation_id": f"file:{relative}",
                "title": path.name,
                "project_id": project_id,
                "project_name": project_name,
                "attribution_basis": "owner_supplied_directory",
                "created_at_epoch": timestamp,
                "updated_at_epoch": timestamp,
                "source_file": relative,
                "source_hash": content_hash,
                "body": body,
                "metadata": {
                    "source_type": "file",
                    "relative_path": relative,
                    "size_bytes": path.stat().st_size,
                    "mime_type": mimetypes.guess_type(path.name)[0],
                    "source_hash": content_hash,
                },
            })
            imported += 1
        except Exception as exc:
            failed += 1
            failures.append({"path": relative, "reason": f"{type(exc).__name__}: {exc}"})

    completed = utc_now()
    run = {
        "run_id": str(uuid.uuid4()),
        "source_name": str(root),
        "source_hash": aggregate.hexdigest(),
        "imported_documents": imported,
        "assigned_documents": imported,
        "unassigned_documents": 0,
        "started_at": started,
        "completed_at": completed,
        "metadata": {
            "source_type": "directory",
            "project_id": project_id,
            "project_name": project_name,
            "skipped": skipped,
            "failed": failed,
            "failures": failures,
        },
    }
    db.record_ingestion(run)
    return run


def ingest_archive(
    db: Database,
    archive_path: Path,
    *,
    source_type: str,
    limits: ArchiveLimits,
    staging_dir: Path | None = None,
    project_id: str | None = None,
    project_name: str | None = None,
) -> dict[str, Any]:
    """Index supported members of a verified ZIP container without extracting it wholesale.

    The bridge intentionally treats an external backup as a container, not as a
    provider-specific schema. Project attribution is retained as UNASSIGNED unless
    the owner explicitly confirms that every indexed member belongs to one project.
    """
    if not archive_path.is_file() or not zipfile.is_zipfile(archive_path):
        raise ValueError("invalid_archive_zip")
    source_type = source_type.strip()
    if not source_type:
        raise ValueError("source_type_required")
    project_id = project_id.strip() if project_id else None
    project_name = project_name.strip() if project_name else None
    if bool(project_id) != bool(project_name):
        raise ValueError("project_id_and_project_name_must_be_supplied_together")
    work_dir = staging_dir or archive_path.parent
    work_dir.mkdir(parents=True, exist_ok=True)
    _require_space(work_dir, 0, limits.min_free_bytes)

    started = utc_now()
    source_hash = sha256_file(archive_path)
    imported = assigned = unassigned = skipped = failed = 0
    reason_counts: dict[str, int] = {}
    samples: list[dict[str, str]] = []
    manifest = hashlib.sha256()

    def record(reason: str, relative: str) -> None:
        reason_counts[reason] = reason_counts.get(reason, 0) + 1
        if len(samples) < MAX_FAILURE_SAMPLES:
            samples.append({"path": relative, "reason": reason})

    if project_id and project_name:
        project = {"project_id": project_id, "project_name": project_name, "basis": "owner_supplied_archive"}
        db.upsert_project(project_id, project_name, project["basis"])
    else:
        project = {
            "project_id": "unassigned",
            "project_name": "UNASSIGNED",
            "basis": "no_explicit_project_metadata",
        }

    with zipfile.ZipFile(archive_path) as archive:
        members = _safe_archive_members(archive, limits)
        for info, relative in members:
            manifest.update(relative.encode("utf-8"))
            manifest.update(b"\0")
            manifest.update(f"{info.CRC:08x}:{info.file_size}:{info.compress_size}".encode("ascii"))
            manifest.update(b"\n")

        with tempfile.TemporaryDirectory(prefix="gpb-member-", dir=work_dir) as temporary:
            staging = Path(temporary)
            for index, (info, relative) in enumerate(members):
                parser = extractor(Path(relative))
                if parser is None:
                    skipped += 1
                    record("unsupported_type", relative)
                    continue
                if info.file_size > limits.max_indexable_file_bytes:
                    skipped += 1
                    record("indexable_file_size_limit", relative)
                    continue

                member_path = staging / f"{index:08d}{Path(relative).suffix.lower()}"
                try:
                    _write_archive_member(archive, info, member_path, limits)
                    text = parser(member_path).strip()
                    if not text:
                        skipped += 1
                        record("empty_extracted_text", relative)
                        continue
                    content_hash = sha256_file(member_path)
                    document_id = "archive:" + hashlib.sha256(
                        (source_hash + "\0" + relative).encode("utf-8")
                    ).hexdigest()
                    body = (
                        f"# {Path(relative).name}\n\n## Provenance\n\n"
                        f"Container: {archive_path.name}\n"
                        f"Container SHA-256: {source_hash}\n"
                        f"Source type: {source_type}\n"
                        f"Project: {project['project_name']}\n"
                        f"Project ID: {project['project_id']}\n"
                        f"Attribution basis: {project['basis']}\n"
                        f"Relative path: {relative}\n"
                        f"Member SHA-256: {content_hash}\n\n## Extracted content\n\n{text}\n"
                    )
                    db.upsert_document({
                        "document_id": document_id,
                        "conversation_id": f"archive:{relative}",
                        "title": Path(relative).name,
                        "project_id": project["project_id"],
                        "project_name": project["project_name"],
                        "attribution_basis": project["basis"],
                        "created_at_epoch": None,
                        "updated_at_epoch": None,
                        "source_file": relative,
                        "source_hash": content_hash,
                        "body": body,
                        "metadata": {
                            "source_type": source_type,
                            "container_format": "zip",
                            "container_name": archive_path.name,
                            "container_sha256": source_hash,
                            "relative_path": relative,
                            "size_bytes": info.file_size,
                            "mime_type": mimetypes.guess_type(relative)[0],
                            "source_hash": content_hash,
                        },
                    })
                    imported += 1
                    if project["project_id"] == "unassigned":
                        unassigned += 1
                    else:
                        assigned += 1
                except (zipfile.BadZipFile, OSError):
                    raise
                except ValueError as exc:
                    if str(exc) == "insufficient_storage":
                        raise
                    failed += 1
                    record(str(exc), relative)
                except Exception as exc:
                    failed += 1
                    record(type(exc).__name__, relative)
                finally:
                    member_path.unlink(missing_ok=True)

    completed = utc_now()
    run = {
        "run_id": str(uuid.uuid4()),
        "source_name": archive_path.name,
        "source_hash": source_hash,
        "imported_documents": imported,
        "assigned_documents": assigned,
        "unassigned_documents": unassigned,
        "started_at": started,
        "completed_at": completed,
        "metadata": {
            "source_type": source_type,
            "container_format": "zip",
            "member_count": len(members),
            "member_manifest_sha256": manifest.hexdigest(),
            "project_id": project_id,
            "project_name": project_name,
            "skipped": skipped,
            "failed": failed,
            "reason_counts": reason_counts,
            "member_samples": samples,
            "member_samples_truncated": sum(reason_counts.values()) > len(samples),
        },
    }
    db.record_ingestion(run)
    return run
